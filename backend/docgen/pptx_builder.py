import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Standard python-pptx default template layout indices: 0 = title slide,
# 1 = "Title and Content" (title placeholder + one body placeholder).
_TITLE_SLIDE_LAYOUT = 0
_CONTENT_SLIDE_LAYOUT = 1
_CHART_POSITION = (Inches(1), Inches(4.5), Inches(8), Inches(2.5))


def build_presentation(structure: dict) -> bytes:
    """Builds a real .pptx from a {"title": str, "slides": [{"heading":
    str, "bullets": [str], "chart": {...} | None}]} structure (the shape
    agents.tasks._run_pptx_generation_step parses from the preceding
    step's raw text). Malformed chart data is skipped, not fatal — the
    deck still ships without that one chart rather than failing the run."""
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_SLIDE_LAYOUT])
    title_slide.shapes.title.text = str(structure.get("title", ""))

    for slide_data in structure.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[_CONTENT_SLIDE_LAYOUT])
        slide.shapes.title.text = str(slide_data.get("heading", ""))

        bullets = slide_data.get("bullets") or []
        body = slide.placeholders[1].text_frame
        if bullets:
            body.text = str(bullets[0])
            for bullet in bullets[1:]:
                body.add_paragraph().text = str(bullet)

        _add_chart_if_valid(slide, slide_data.get("chart"))

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_chart_if_valid(slide, chart_data: dict | None) -> None:
    if not chart_data:
        return
    try:
        categories = [str(c) for c in chart_data["categories"]]
        values = [float(v) for v in chart_data["values"]]
    except (KeyError, TypeError, ValueError):
        return
    if not categories or not values or len(categories) != len(values):
        return

    data = CategoryChartData()
    data.categories = categories
    data.add_series(str(chart_data.get("chart_title", "")), values)
    x, y, cx, cy = _CHART_POSITION
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, data)
